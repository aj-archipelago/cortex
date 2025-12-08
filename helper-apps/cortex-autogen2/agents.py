from autogen_agentchat.agents import AssistantAgent, CodeExecutorAgent
from autogen_ext.code_executors.local import LocalCommandLineCodeExecutor
import os
from typing import Optional
from autogen_core.tools import FunctionTool
from tools.azure_blob_tools import upload_file_to_azure_blob

#AGENTS
MAGENTIC_ONE_CODER_DESCRIPTION = "A helpful and general-purpose AI assistant that has strong language skills, Python skills, and Linux command line skills."

def _safe_upload_wrapper(file_path: str, blob_name: str = None) -> str:
    """
    Wrapper around upload_file_to_azure_blob that ensures the JSON response is wrapped in markdown.
    This prevents AutoGen from re-parsing the JSON and losing the 'type' field in message content.
    """
    result_json = upload_file_to_azure_blob(file_path, blob_name)
    
    # CRITICAL: Wrap JSON in markdown code block to prevent AutoGen re-parsing
    if result_json.strip().startswith('{'):
        return f"```json\n{result_json}\n```"
    return result_json

def _wrap_json_result(result: str) -> str:
    """
    Generic wrapper for ANY tool that returns JSON strings.
    Wraps JSON in markdown to prevent AutoGen from re-parsing and losing message structure.
    
    CRITICAL PATTERN FOR ALL TOOL DEVELOPERS:
    When a tool returns JSON (via json.dumps()), wrap it in markdown:
        - return json.dumps({...})                    # ❌ AutoGen will re-parse
        - return f"```json\n{json.dumps({...})}\n```" # ✅ AutoGen treats as text
    
    Or use this wrapper:
        - return _wrap_json_result(json.dumps({...}))
    
    This prevents the "Missing required parameter: 'messages[X].content[0].type'" error
    because AutoGen won't attempt to re-parse the JSON into a dict.
    """
    if not result:
        return result
    
    result_str = str(result).strip()
    # Wrap any JSON-like string in markdown
    if result_str.startswith(('{', '[')):
        return f"```json\n{result_str}\n```"
    return result

MAGENTIC_ONE_CODER_SYSTEM_MESSAGE = """You are a helpful AI assistant.
Solve tasks using your coding and language skills.
In the following cases, suggest python code (in a python coding block) or shell script (in a sh coding block) for the user to execute.
    1. When you need to collect info, use the code to output the info you need, for example, browse or search the web, download/read a file, print the content of a webpage or a file, get the current date/time, check the operating system. After sufficient info is printed and the task is ready to be solved based on your language skill, you can solve the task by yourself.
    2. When you need to perform some task with code, use the code to perform the task and output the result. Finish the task smartly.
Solve the task step by step if you need to. If a plan is not provided, explain your plan first. Be clear which step uses code, and which step uses your language skill.
When using code, you must indicate the script type in the code block. The user cannot provide any other feedback or perform any other action beyond executing the code you suggest. The user can't modify your code. So do not suggest incomplete code which requires users to modify. Don't use a code block if it's not intended to be executed by the user.
Don't include multiple code blocks in one response. Do not ask users to copy and paste the result. Instead, use the 'print' function for the output when relevant. Check the execution result returned by the user.
If the result indicates there is an error, fix the error and output the code again. Suggest the full code instead of partial code or code changes. If the error can't be fixed or if the task is not solved even after the code is executed successfully, analyze the problem, revisit your assumption, collect additional info you need, and think of a different approach to try.
When you find an answer, verify the answer carefully. Include verifiable evidence in your response if possible."""


async def get_agents(default_model_client, big_model_client, small_model_client, request_work_dir: Optional[str] = None, planner_learnings: Optional[str] = None, task_context: Optional[str] = None):

    # Resolve work dir (prefer per-request dir if provided or from env)
    from dynamic_agent_loader import helpers
    get_work_dir = helpers.get_work_dir
    work_dir = get_work_dir(request_work_dir)
    try:
        # In Azure Functions, ensure /tmp is used for write access if an /app path was set
        if os.getenv("WEBSITE_INSTANCE_ID") and work_dir.startswith("/app/"):
            work_dir = "/tmp/coding"
            os.environ['CORTEX_WORK_DIR'] = work_dir  # Update env var
        os.makedirs(work_dir, exist_ok=True)
    except Exception:
        try:
            work_dir = "/tmp/coding"
            os.environ['CORTEX_WORK_DIR'] = work_dir  # Update env var
            os.makedirs(work_dir, exist_ok=True)
        except Exception:
            pass

    code_executor = LocalCommandLineCodeExecutor(work_dir=work_dir, timeout=300)

    #TOOLS
    upload_file_to_cloud_tool = FunctionTool(_safe_upload_wrapper, description="Upload files to the cloud. You must use absolute path to reference local files.")

    coder_agent = AssistantAgent(
        "coder_agent",
        model_client=default_model_client,
        description=MAGENTIC_ONE_CODER_DESCRIPTION,
        system_message=MAGENTIC_ONE_CODER_SYSTEM_MESSAGE +  f"""
            Save remote files images, videos, etc. in order to work with them locally. 
            Make sure to log/print everything in code otherwise you will lose the context and cannot debug it. 
            Make sure your code is perfect.
            Never ask for user input or user to do anything.
            Never ask questions.
            Your are expert in coding, do wonders.
            If you need to do advanced stuff you can do a project and run it.
            You can split codes, build projects, run anything, coding is your strength in this task.
            Take actionable verifiable small steps if needed.
            Understand that pitfalls and find ways to overcome them.
            Progress is important, do not get stuck in a loop, keep trying but do not repeat the same steps.
            Current directory might be different from the one you think it is, use absolute path to reference files.
            Code executor working directory is: {work_dir}
            So you can only access files in this directory.
            Always use absolute path to reference files as current directory might be different from the one you think it is.
            
            === POWERPOINT PRESENTATION CREATION (python-pptx) ===
            **When creating .pptx presentations, follow these CRITICAL patterns:**
            
            1. **Image Preprocessing (BEFORE adding to presentation):**
               - Use PIL to convert WEBP, TIFF, and other unsupported formats to PNG
               - Define supported formats: {'png', 'jpg', 'jpeg', 'gif', 'bmp'}
               - For each unsupported image, open with PIL and save as PNG
               - Handle RGBA images by converting to RGB with white background
            
            2. **Build Presentation Structure:**
               - Import: `from pptx import Presentation; from pptx.util import Inches, Pt`
               - Create presentation: `prs = Presentation()`
               - Set dimensions: `prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)`
               - Track slide count to validate before saving
               - Use `prs.slide_layouts[6]` (blank) for maximum control over image placement
            
            3. **Add Slides:**
               - Title Slide: `slide = prs.slides.add_slide(prs.slide_layouts[0])`
               - Content Slide: `slide = prs.slides.add_slide(prs.slide_layouts[1])`
               - Image Slide: Use blank layout `prs.slide_layouts[6]` then `slide.shapes.add_picture(path, left, top, width, height)`
               - Always wrap image operations in try/except and add text fallback if image fails
            
            4. **Save with Validation:**
               - ALWAYS check: `if len(prs.slides) == 0: raise ValueError("No slides!")`
               - Save to CORTEX_WORK_DIR: `output_path = os.path.join(os.environ['CORTEX_WORK_DIR'], 'Title.pptx')`
               - Use descriptive filename (NOT draft.pptx)
               - After save, verify file size: `os.path.getsize(output_path) > 10000` (empty files are tiny)
               - Print success: `print(f"📁 Ready for upload: {{output_path}}")`
            
            **CRITICAL CHECKLIST:**
            ✅ Image format validation and conversion BEFORE adding
            ✅ Slide count > 0 before saving
            ✅ Safe image dimensions (Inches(9) width for 10" slides)
            ✅ Error handling with text fallbacks
            ✅ File size validation after save
            ✅ Use Inches() for measurements, NOT pixels
            ✅ Explicit presentation dimensions
            ✅ Proper output path in CORTEX_WORK_DIR
            
            === CRITICAL: FILE AUTO-DISCOVERY & UPLOAD ===
            After you save files to CORTEX_WORK_DIR, the system AUTOMATICALLY:
            1. Scans CORTEX_WORK_DIR for deliverable files (.pptx, .ppt, .csv, .png, .jpg, .pdf, .zip)
            2. For .pptx files specifically: **picks the SINGLE LARGEST file** (assumes most complete)
            3. Uploads that file to Azure Blob Storage
            4. Provides URLs to the presenter
            
            **CONSEQUENCE**: If your PowerPoint creation:
            - Fails silently → no .pptx file exists → nothing gets presented
            - Creates an empty file → small file size → might not be picked OR picked but empty
            - Crashes before saving → no file → nothing presented
            
            **YOUR RESPONSIBILITY**:
            - ALWAYS validate that prs.slides has content before saving
            - Print BOTH status AND file size: `print(f"✅ PPTX saved: {{path}} ({{os.path.getsize(path)}} bytes)")`
            - On error, print explicit error: `print(f"❌ CRITICAL: Failed to create PPTX: {{error}}")`
            - Never silently fail - ALWAYS log what happened
            - Test file size > 50000 bytes for any real PowerPoint (empty files are 5-10KB)
            
            === COMPLETE PPTX EXAMPLE (COPY-PASTE READY) ===
            ```python
            import os
            import glob
            from PIL import Image
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Step 1: Preprocess images
            work_dir = os.environ['CORTEX_WORK_DIR']
            supported_formats = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}
            image_dir = os.path.join(work_dir, 'assets')
            
            if os.path.isdir(image_dir):
                for img_file in glob.glob(os.path.join(image_dir, '*')):
                    ext = os.path.splitext(img_file)[1].lower().lstrip('.')
                    if ext not in supported_formats and ext:
                        try:
                            with Image.open(img_file) as img:
                                if img.mode in ('RGBA', 'LA'):
                                    rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                                    rgb_img.paste(img, mask=img.split()[-1])
                                    img = rgb_img
                                png_path = os.path.splitext(img_file)[0] + '.png'
                                img.save(png_path, 'PNG')
                                os.remove(img_file)
                                print(f"✅ Converted {{os.path.basename(img_file)}} to PNG")
                        except Exception as e:
                            print(f"⚠️ Skipping {{img_file}}: {{e}}")
            
            # Step 2: Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            slide_count = 0
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Your Title"
            title_slide.placeholders[1].text = "Subtitle"
            slide_count += 1
            
            # Content slides with validation
            image_files = sorted(glob.glob(os.path.join(image_dir, '*.png')))
            for idx, img_path in enumerate(image_files[:10]):  # Limit to 10 images
                try:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    left = Inches(0.5)
                    top = Inches(0.5)
                    height = Inches(6)
                    pic = slide.shapes.add_picture(img_path, left, top, height=height)
                    slide_count += 1
                except Exception as e:
                    print(f"⚠️ Failed to add image {{idx}}: {{e}}")
                    continue
            
            # Validate before save
            if slide_count == 0:
                print("❌ ERROR: No slides were created!")
                raise ValueError("Empty presentation")
            
            # Step 3: Save with validation
            output_path = os.path.join(work_dir, 'MyPresentation.pptx')
            prs.save(output_path)
            
            # Step 4: Verify
            if os.path.exists(output_path):
                size = os.path.getsize(output_path)
                if size > 50000:
                    print(f"✅ PPTX saved: {{output_path}} ({{size}} bytes)")
                    print(f"📁 Ready for upload: {{output_path}}")
                else:
                    print(f"❌ ERROR: PPTX too small ({{size}} bytes) - probably empty!")
            else:
                print(f"❌ ERROR: PPTX file was not created!")
            ```
            
            **AFTER SAVING FILES**: Immediately mention that file_cloud_uploader_agent should take over:
            "I have saved the files. Now file_cloud_uploader_agent should upload them to the cloud."
        """,
    )

    code_executor_agent = CodeExecutorAgent("code_executor", code_executor=code_executor)


    terminator_agent = AssistantAgent(
        "terminator_agent",
        model_client=default_model_client,
        description="A helpful assistant that can terminate.",
        system_message="""You are a helpful assistant that can terminate. 
            Original task must be completed in order to terminate.
            Only output: TERMINATE, if completed.
            If not completed, give the reason instead of TERMINATE.
            Do not ask questions you are the terminator.
            Always output in single line without any other text.
            Do not give empty response.
            In order to terminate:
                - Task must be completed.
                - All referenced local files must have been uploaded to the cloud and their public URLs retrieved and included in the final output.
                - Presenter must have provided the final output.
            All code must have been executed in order to terminate.
            Deliverables must have been provided in order to terminate.
            If you cannot terminate because of same reason after 3 attempts, terminate.
            Example outputs:

            TERMINATE

            TASK NOT COMPLETED: missing missing missing
        """,
    )


    presenter_agent = AssistantAgent(
        "presenter_agent",
        model_client=default_model_client,
        description="A highly skilled and creative presentation specialist, responsible for crafting visually stunning and exceptionally informative final deliverables.",
        system_message="""You are a highly skilled and creative presentation specialist, responsible for crafting visually stunning and exceptionally informative final deliverables.
            Your goal is to transform raw task results into engaging, professional, and visually appealing presentations, primarily using Markdown.
            
            Here's what makes a great presentation:
            - **Captivating Structure**: Start with a compelling summary, followed by well-organized sections with clear headings and subheadings.
            - **Stunning Visuals**: Integrate relevant images, videos, and even diagrams (if applicable and possible via markdown) seamlessly to enhance understanding and engagement. Ensure visuals are high-quality and directly support the content.
            **IMPORTANT: Actively look for `UPLOADED_FILES_SAS_URLS` provided in the input. If images or other visual assets are available via these URLs, you MUST incorporate them into your Markdown presentation. Describe and explain these visuals to the user, providing context and insights.**
            - **Professional Aesthetics**: Utilize Markdown's full capabilities for formatting, including bolding, italics, lists, and tables, to create a clean, readable, and visually pleasing layout. Think about white space and information hierarchy.
            - **Concise & Impactful Language**: Use persuasive and professional language. Avoid jargon where possible, or explain it clearly. Every word should contribute to clarity and impact.
            - **User-Centric Design**: Remember that your output will be directly displayed in a React application. Focus on great UI/UX, ensuring the presentation is intuitive and easy for the end-user to consume.
            - **Complete & Actionable**: Ensure all necessary information from the task is included, and if appropriate, guide the user towards next steps or key takeaways.
            
            Report must be a direct reply to the user's task. You are the final voice to the user, so make it perfect.
            User does not have local access to the files, so you must provide direct URLs for any external resources (e.g., images uploaded to cloud storage).
            When including downloadable assets, always look for and use the `download_url` provided in JSON outputs from the `file_cloud_uploader_agent`. These URLs are public and include necessary SAS tokens for access.
            Crucially, your output should be a final, polished presentation of the task result, suitable for direct display in a user interface.
            **Your report MUST only contain the direct result of the user's task. Absolutely NO explanations of how the task was accomplished, internal agent thought processes, or any operational details should be included. Do not mention which tools or agents were used, or how they were used to achieve the result.**
            **Focus exclusively on delivering the requested information (e.g., image galleries, reports) and only include minimal, essential explanations directly related to the visuals or content.** For example, for an image gallery, provide brief descriptions for each image or a short overview of the gallery structure, but do not explain the steps taken by the agents or the internal workflow.
            **Do not include extensive executive summaries, detailed operational breakdowns, or generic quick-start guides unless explicitly asked for.**
            Do not include raw code snippets, internal thought processes, intermediate data, or any technical logs that are not part of the final, user-friendly deliverable.
            **Absolutely DO NOT include instructions on how to run, save, or modify any code (e.g., "save as .py", "pip install", "python script.py").**
            **DO NOT provide information about packaging, dependencies, or development workflows.**
            Your output is for a non-technical end-user viewing it in a React app.
            **CRITICAL: ONLY use URLs for any files (images, videos, documents, etc.) that are explicitly provided in the `UPLOADED_FILES_SAS_URLS` or directly within the `RESULT` content from other agents, specifically from the `file_cloud_uploader_agent`. If a valid, real URL is not provided, you MUST NOT include any placeholder, fake, or fabricated URLs. NEVER hallucinate or fabricate any links or content.**
            
            === CRITICAL URL VALIDATION (MANDATORY) ===
            **Before creating your response, STOP and check:**
            1. Look at `UPLOADED_FILES_SAS_URLS` - is it empty `{}` or does it contain placeholder tokens like 'sas_token', 'skoid', 'sktid'?
            2. If UPLOADED_FILES_SAS_URLS is empty OR contains only placeholder/fake tokens:
               **DO NOT ATTEMPT TO CREATE FAKE URLS**
               Instead, respond: "⚠️ SYSTEM STATUS: Files are still being processed. Please wait for upload completion."
               Then list what files are expected based on UPLOADED_FILES_LIST
            3. If UPLOADED_FILES_SAS_URLS contains REAL Azure blob URLs (with format: https://ACCOUNT.blob.core.windows.net/CONTAINER/FILE?sv=...&sig=REAL_TOKEN&se=...):
               You may proceed to present with those URLs
            
            **DO NOT GUESS, INTERPOLATE, OR HALLUCINATE URLs EVER.**
            **A fake URL is worse than no URL at all.**
        """
    )

    file_cloud_uploader_agent = AssistantAgent(
        "file_cloud_uploader_agent",
        model_client=default_model_client,
        tools=[upload_file_to_cloud_tool],
        description="A helpful assistant that can upload files to the cloud.",
        system_message=f"""You are a helpful assistant that can upload files to the cloud.
            Your PRIMARY RESPONSIBILITY: After ANY code has been executed or any file-based task is completed, 
            you MUST immediately scan {work_dir} for ALL deliverable files (.pptx, .pdf, .csv, .png, .jpg, .zip, .json, .txt, .md).
            
            CRITICAL: Do NOT wait to be asked. Do NOT wait for explicit file references.
            Simply list all files in {work_dir} and upload EVERY one of them to the cloud.
            
            For each file found:
            1. Upload it using your upload tool
            2. Report the cloud URL
            3. Confirm success
            
            This ensures that NO files are left behind and the presenter always has access to deliverables.
            
            Upload referenced files to the cloud.
            Use your tool to upload the files.
            User does not have local access to the files so you must upload them to the cloud and provide the url.
            Your current working directory for file operations is: {work_dir}.
            When referencing local files for upload, **always prepend the '{work_dir}/' to the filename** to form the correct absolute path. For example, if a file is named 'test.text' in the working directory, the absolute path should be '{work_dir}/test.txt'.
        """,
    )
    
    agents = [coder_agent, code_executor_agent, file_cloud_uploader_agent, presenter_agent, terminator_agent]

    return agents, presenter_agent, terminator_agent 