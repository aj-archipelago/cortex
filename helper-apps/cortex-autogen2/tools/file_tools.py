"""
Universal File Tools for Cortex-AutoGen2
Enhanced file handling system that works across all file types.
"""

import asyncio
import logging
import os
import json
import mimetypes
from pathlib import Path
from typing import List, Optional, Dict, Any, Union
from autogen_core.tools import FunctionTool

logger = logging.getLogger(__name__)


class UniversalFileHandler:
    """Universal file handler that intelligently processes any file type."""
    
    @staticmethod
    def detect_file_type(file_path: str) -> Dict[str, Any]:
        """
        Intelligently detect file type and characteristics.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary with file type information
        """
        if not os.path.exists(file_path):
            return {"error": "File not found", "exists": False}
        
        filename = os.path.basename(file_path)
        extension = os.path.splitext(filename)[1].lower()
        size = os.path.getsize(file_path)
        
        # MIME type detection
        mime_type, _ = mimetypes.guess_type(file_path)
        
        # Category classification
        categories = {
            "document": [".pdf", ".doc", ".docx", ".txt", ".md", ".rtf"],
            "spreadsheet": [".xls", ".xlsx", ".csv", ".tsv"],
            "presentation": [".ppt", ".pptx"],
            "image": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp"],
            "data": [".json", ".xml", ".yaml", ".yml"],
            "code": [".py", ".js", ".html", ".css", ".sql", ".r"],
            "archive": [".zip", ".tar", ".gz", ".rar", ".7z"],
            "media": [".mp4", ".mp3", ".avi", ".mov", ".wav"]
        }
        
        category = "unknown"
        for cat, extensions in categories.items():
            if extension in extensions:
                category = cat
                break
        
        return {
            "filename": filename,
            "extension": extension,
            "size": size,
            "mime_type": mime_type,
            "category": category,
            "exists": True,
            "is_deliverable": not filename.startswith(("tmp_", "temp_", ".")),
            "is_readable": category in ["document", "data", "code"] or extension in [".txt", ".csv", ".json"],
            "description": _get_file_description(extension, category)
        }
    
    @staticmethod
    def read_file_intelligently(file_path: str, max_length: int = 5000) -> Dict[str, Any]:
        """
        Intelligently read file content based on file type.
        
        Args:
            file_path: Path to the file
            max_length: Maximum characters to read for text files
            
        Returns:
            Dictionary with file content and metadata
        """
        file_info = UniversalFileHandler.detect_file_type(file_path)
        
        if not file_info["exists"]:
            return file_info
        
        try:
            content_info = {
                **file_info,
                "content": None,
                "preview": None,
                "metadata": {}
            }
            
            if file_info["category"] == "image":
                content_info["preview"] = f"Image file: {file_info['filename']} ({file_info['size']} bytes)"
                content_info["metadata"]["viewable"] = True
                
            elif file_info["category"] == "document" and file_info["extension"] == ".pdf":
                content_info["preview"] = f"PDF document: {file_info['filename']} ({file_info['size']} bytes)"
                content_info["metadata"]["pages"] = "Unknown"
                content_info["metadata"]["viewable"] = True
                
            elif file_info["is_readable"]:
                # Try to read as text
                encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
                content = None
                
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read(max_length)
                        break
                    except (UnicodeDecodeError, UnicodeError):
                        continue
                
                if content is not None:
                    content_info["content"] = content
                    content_info["preview"] = content[:500] + "..." if len(content) > 500 else content
                    content_info["metadata"]["encoding"] = encoding
                    content_info["metadata"]["truncated"] = len(content) >= max_length
                    
                    # Special handling for structured data
                    if file_info["extension"] == ".json":
                        try:
                            json_data = json.loads(content)
                            content_info["metadata"]["json_valid"] = True
                            content_info["metadata"]["json_type"] = type(json_data).__name__
                            if isinstance(json_data, list):
                                content_info["metadata"]["json_length"] = len(json_data)
                        except json.JSONDecodeError:
                            content_info["metadata"]["json_valid"] = False
                            
                    elif file_info["extension"] == ".csv":
                        lines = content.split('\n')
                        content_info["metadata"]["csv_rows"] = len(lines)
                        content_info["metadata"]["csv_columns"] = len(lines[0].split(',')) if lines else 0
                else:
                    content_info["preview"] = "Binary file - cannot preview as text"
            else:
                content_info["preview"] = f"Binary file: {file_info['filename']} ({file_info['size']} bytes)"
                content_info["metadata"]["binary"] = True
            
            return content_info
            
        except Exception as e:
            return {
                **file_info,
                "error": f"Error reading file: {str(e)}",
                "content": None,
                "preview": None
            }


def _get_file_description(extension: str, category: str) -> str:
    """Get a human-readable description of the file type."""
    descriptions = {
        ".pdf": "PDF Document",
        ".txt": "Text File",
        ".csv": "CSV Data File", 
        ".json": "JSON Data File",
        ".png": "PNG Image",
        ".jpg": "JPEG Image",
        ".jpeg": "JPEG Image",
        ".gif": "GIF Image",
        ".svg": "SVG Vector Image",
        ".html": "HTML Document",
        ".xml": "XML Document",
        ".docx": "Word Document",
        ".xlsx": "Excel Spreadsheet",
        ".pptx": "PowerPoint Presentation",
        ".zip": "ZIP Archive",
        ".py": "Python Script",
        ".js": "JavaScript File",
        ".css": "CSS Stylesheet",
        ".md": "Markdown Document",
        ".sql": "SQL Script",
        ".yaml": "YAML Configuration",
        ".yml": "YAML Configuration"
    }
    
    return descriptions.get(extension, f"{category.title()} File" if category != "unknown" else "Unknown File Type")


async def download_image(url: str, filename: str, work_dir: Optional[str] = None) -> str:
    """
    Downloads an image from a URL and saves it to the working directory.

    Args:
        url: The URL of the image to download.
        filename: The local filename to save the image as.
        work_dir: Optional working directory path.

    Returns:
        A JSON string indicating success or failure.
    """
    if not work_dir:
        work_dir = os.getcwd()
    
    file_path = os.path.join(work_dir, filename)
    
    try:
        import requests
        BROWSER_UA = (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/125.0.0.0 Safari/537.36"
        )
        session = requests.Session()
        session.headers.update({
            "User-Agent": BROWSER_UA,
            "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
            "Referer": "https://www.google.com/",
            "Cache-Control": "no-cache",
        })

        # Attempt to derive an original Wikimedia URL if this is a thumbnail
        wm_orig = None
        try:
            if "upload.wikimedia.org" in url and "/thumb/" in url:
                parts = url.split("/thumb/")
                if len(parts) == 2:
                    tail = parts[1]
                    segs = tail.split("/")
                    if len(segs) >= 3:
                        wm_orig = parts[0] + "/" + segs[0] + "/" + segs[1] + "/" + segs[2]
        except Exception:
            wm_orig = None

        candidates = []
        if wm_orig:
            candidates.append(wm_orig)
        candidates.append(url)

        last_err = None
        for candidate in candidates:
            try:
                with session.get(candidate, stream=True, timeout=25, allow_redirects=True) as response:
                    response.raise_for_status()

                    content_type = (response.headers.get("Content-Type") or "").lower()

                    # Peek first few bytes to validate image magic if header is missing/misleading
                    first_chunk = next(response.iter_content(chunk_size=4096), b"")

                    def looks_like_image(buf: bytes) -> bool:
                        if not buf or len(buf) < 4:
                            return False
                        sigs = [
                            b"\x89PNG\r\n\x1a\n",  # PNG
                            b"\xff\xd8\xff",        # JPEG
                            b"GIF87a", b"GIF89a",    # GIF
                            b"RIFF"                    # WEBP starts with RIFF
                        ]
                        return any(buf.startswith(sig) for sig in sigs)

                    if not (content_type.startswith("image/") or looks_like_image(first_chunk)):
                        last_err = f"Non-image content-type: {content_type} for {candidate}"
                        continue

                    # Write first chunk then stream the rest
                    with open(file_path, 'wb') as f:
                        if first_chunk:
                            f.write(first_chunk)
                        for chunk in response.iter_content(chunk_size=8192):
                            if chunk:
                                f.write(chunk)

                logger.info(f"✅ Successfully downloaded image from {candidate} to {file_path}")
                return json.dumps({"status": "success", "file_path": file_path})
            except Exception as e:
                last_err = str(e)
                continue

        logger.error(f"❌ Failed to download image after candidates. Last error: {last_err}")
        return json.dumps({"status": "error", "message": last_err or "download_failed"})
    except Exception as e:
        logger.error(f"❌ Failed to download image from {url}: {e}")
        return json.dumps({"status": "error", "message": str(e)})


# Enhanced file tools
async def list_files_in_work_dir(work_dir: Optional[str] = None) -> str:
    """
    Intelligently list and categorize all files in the working directory.
    
    Args:
        work_dir: Optional working directory path
    
    Returns:
        String containing categorized file listing with metadata
    """
    if not work_dir or not os.path.exists(work_dir):
        return "❌ Working directory not found or not specified"
    
    try:
        all_files = []
        categories = {
            "deliverable": [],
            "temporary": [],
            "data": [],
            "documents": [],
            "images": [],
            "code": [],
            "other": []
        }
        
        # CRITICAL: Search recursively to find ALL files in subdirectories
        for root, dirs, files in os.walk(work_dir):
            # Skip hidden directories and common non-deliverable directories
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in ['__pycache__', 'node_modules']]
            
            for item in files:
                item_path = os.path.join(root, item)
                if os.path.isfile(item_path):
                    file_info = UniversalFileHandler.detect_file_type(item_path)
                    # Add relative path from work_dir for better readability
                    file_info["relative_path"] = os.path.relpath(item_path, work_dir)
                    all_files.append(file_info)
                    
                    # Categorize for display
                    if not file_info["is_deliverable"]:
                        categories["temporary"].append(file_info)
                    elif file_info["category"] == "document":
                        categories["documents"].append(file_info)
                    elif file_info["category"] == "image":
                        categories["images"].append(file_info)
                    elif file_info["category"] == "data":
                        categories["data"].append(file_info)
                    elif file_info["category"] == "code":
                        categories["code"].append(file_info)
                    else:
                        categories["deliverable"].append(file_info)
        
        if not all_files:
            return f"📁 No files found in working directory: {work_dir}"
        
        # Sort each category by size (largest first)
        for category in categories.values():
            category.sort(key=lambda x: x["size"], reverse=True)
        
        result = f"📁 **UNIVERSAL FILE DISCOVERY**\n"
        result += f"Directory: {work_dir}\n"
        result += f"Total Files: {len(all_files)}\n\n"
        
        # Show deliverable files first
        deliverable_count = sum(len(cat) for cat_name, cat in categories.items() if cat_name != "temporary")
        
        if deliverable_count > 0:
            result += f"🎯 **DELIVERABLE FILES** ({deliverable_count} files):\n\n"
            
            if categories["documents"]:
                result += f"📄 **Documents** ({len(categories['documents'])}):\n"
                for file_info in categories["documents"]:
                    # Use absolute path stored in file_info
                    abs_path = file_info.get('absolute_path', os.path.join(work_dir, file_info.get('relative_path', file_info['filename'])))
                    file_path = file_info.get('relative_path', file_info['filename'])
                    result += f"  ✅ {file_path} ({file_info['size']} bytes) - {file_info['description']}\n"
                    result += f"     Path: {abs_path}\n"
                result += "\n"
            
            if categories["images"]:
                result += f"🖼️ **Images** ({len(categories['images'])}):\n"
                for file_info in categories["images"]:
                    # Use absolute path stored in file_info
                    abs_path = file_info.get('absolute_path', os.path.join(work_dir, file_info.get('relative_path', file_info['filename'])))
                    file_path = file_info.get('relative_path', file_info['filename'])
                    result += f"  ✅ {file_path} ({file_info['size']} bytes) - {file_info['description']}\n"
                    result += f"     Path: {abs_path}\n"
                result += "\n"
            
            if categories["data"]:
                result += f"📊 **Data Files** ({len(categories['data'])}):\n"
                for file_info in categories["data"]:
                    # Use absolute path stored in file_info
                    abs_path = file_info.get('absolute_path', os.path.join(work_dir, file_info.get('relative_path', file_info['filename'])))
                    file_path = file_info.get('relative_path', file_info['filename'])
                    result += f"  ✅ {file_path} ({file_info['size']} bytes) - {file_info['description']}\n"
                    result += f"     Path: {abs_path}\n"
                result += "\n"
            
            if categories["deliverable"]:
                result += f"📦 **Other Deliverables** ({len(categories['deliverable'])}):\n"
                for file_info in categories["deliverable"]:
                    # Use absolute path stored in file_info
                    abs_path = file_info.get('absolute_path', os.path.join(work_dir, file_info.get('relative_path', file_info['filename'])))
                    file_path = file_info.get('relative_path', file_info['filename'])
                    result += f"  ✅ {file_path} ({file_info['size']} bytes) - {file_info['description']}\n"
                    result += f"     Path: {abs_path}\n"
                result += "\n"
        
        if categories["temporary"]:
            result += f"🗂️ **Temporary Files** ({len(categories['temporary'])}):\n"
            for file_info in categories["temporary"]:
                result += f"  - {file_info['filename']} ({file_info['size']} bytes)\n"
            result += "\n"
        
        if categories["code"]:
            result += f"💻 **Code Files** ({len(categories['code'])}):\n"
            for file_info in categories["code"]:
                result += f"  📝 {file_info['filename']} ({file_info['size']} bytes) - {file_info['description']}\n"
        
        # Add recommendations with absolute paths list
        if deliverable_count > 0:
            result += f"\n🚀 **NEXT ACTIONS**:\n"
            result += f"1. Read and analyze deliverable files\n"
            result += f"2. Upload all deliverables to Azure\n"
            result += f"3. Provide download links to user\n\n"
            result += "📋 **ABSOLUTE PATHS FOR UPLOAD** (copy these paths for batch_upload_for_uploader):\n"
            all_deliverable_paths = []
            for category_name in ["documents", "images", "data", "deliverable"]:
                for file_info in categories[category_name]:
                    abs_path = file_info.get('absolute_path', os.path.join(work_dir, file_info.get('relative_path', file_info['filename'])))
                    all_deliverable_paths.append(abs_path)
            for path in all_deliverable_paths:
                result += f"  - {path}\n"
        else:
            result += f"\n⚠️ **WARNING**: No deliverable files found!\n"
            result += f"Check if the task was completed successfully.\n"
        
        return result
        
    except Exception as e:
        return f"❌ Error listing files: {str(e)}"

async def create_file(filename: str, content: str, work_dir: Optional[str] = None) -> str:
    """
    Creates a new file with the given content in the specified working directory.

    Args:
        filename: The name of the file to create.
        content: The content to write to the file.
        work_dir: Optional working directory path. If not provided, uses the current directory.

    Returns:
        A JSON string confirming the file creation or an error message.
    """
    dir_path = work_dir or os.getcwd()
    file_path = os.path.join(dir_path, filename)

    try:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)
        return json.dumps({"status": "success", "message": f"File '{filename}' created successfully in '{dir_path}'."})
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})

async def read_file_from_work_dir(filename: str, work_dir: Optional[str] = None, max_length: int = 5000) -> str:
    """
    Intelligently read and analyze any file type.
    
    Args:
        filename: Name of the file to read
        work_dir: Optional working directory path
        max_length: Maximum characters to read for text files
    
    Returns:
        String containing file analysis and content
    """
    if not work_dir or not os.path.exists(work_dir):
        return "❌ Working directory not found or not specified"
    
    file_path = os.path.join(work_dir, filename)
    
    if not os.path.exists(file_path):
        return f"❌ File not found: {filename}"
    
    try:
        content_info = UniversalFileHandler.read_file_intelligently(file_path, max_length)
        
        if "error" in content_info:
            return f"❌ {content_info['error']}"
        
        result = f"📄 **UNIVERSAL FILE ANALYSIS: {filename}**\n"
        result += f"Type: {content_info['description']}\n"
        result += f"Category: {content_info['category'].title()}\n"
        result += f"Size: {content_info['size']} bytes\n"
        result += f"MIME Type: {content_info['mime_type'] or 'Unknown'}\n"
        result += f"Deliverable: {'✅ Yes' if content_info['is_deliverable'] else '❌ No'}\n"
        
        if content_info["metadata"]:
            result += f"\n📊 **Metadata**:\n"
            for key, value in content_info["metadata"].items():
                result += f"  {key}: {value}\n"
        
        if content_info["preview"]:
            result += f"\n📝 **Content Preview**:\n"
            if content_info["content"] and content_info["is_readable"]:
                result += f"```\n{content_info['preview']}\n```\n"
                
                if content_info["metadata"].get("truncated"):
                    result += f"\n📏 **Note**: Content truncated at {max_length} characters\n"
            else:
                result += f"{content_info['preview']}\n"
        
        # Add recommendations
        result += f"\n💡 **Recommendations**:\n"
        if content_info["is_deliverable"]:
            result += f"✅ Upload this file to Azure for user download\n"
            if content_info["category"] == "document":
                result += f"📋 This appears to be a document - perfect for user delivery\n"
            elif content_info["category"] == "image":
                result += f"🖼️ This is an image file - should be viewable after upload\n"
            elif content_info["category"] == "data":
                result += f"📊 This contains data - useful for analysis or download\n"
        else:
            result += f"🗂️ This appears to be a temporary file\n"
        
        return result
        
    except Exception as e:
        return f"❌ Error reading file {filename}: {str(e)}"


async def get_file_info(filename: str, work_dir: Optional[str] = None) -> str:
    """
    Get comprehensive metadata about any file type.
    
    Args:
        filename: Name of the file to analyze
        work_dir: Optional working directory path
    
    Returns:
        String containing detailed file metadata
    """
    if not work_dir or not os.path.exists(work_dir):
        return "❌ Working directory not found or not specified"
    
    file_path = os.path.join(work_dir, filename)
    
    if not os.path.exists(file_path):
        return f"❌ File not found: {filename}"
    
    try:
        import time
        stat = os.stat(file_path)
        file_info = UniversalFileHandler.detect_file_type(file_path)
        
        result = f"📊 **COMPREHENSIVE FILE METADATA: {filename}**\n"
        result += f"Full Path: {file_path}\n"
        result += f"Size: {stat.st_size} bytes ({stat.st_size / 1024:.1f} KB)\n"
        result += f"Created: {time.ctime(stat.st_ctime)}\n"
        result += f"Modified: {time.ctime(stat.st_mtime)}\n"
        result += f"Extension: {file_info['extension'] or 'None'}\n"
        result += f"MIME Type: {file_info['mime_type'] or 'Unknown'}\n"
        result += f"Category: {file_info['category'].title()}\n"
        result += f"Description: {file_info['description']}\n"
        result += f"Readable: {'✅ Yes' if file_info['is_readable'] else '❌ No'}\n"
        result += f"Deliverable: {'✅ Yes' if file_info['is_deliverable'] else '❌ No'}\n"
        
        # File permissions
        result += f"\n🔐 **Permissions**:\n"
        result += f"Readable: {'✅' if os.access(file_path, os.R_OK) else '❌'}\n"
        result += f"Writable: {'✅' if os.access(file_path, os.W_OK) else '❌'}\n"
        result += f"Executable: {'✅' if os.access(file_path, os.X_OK) else '❌'}\n"
        
        return result
        
    except Exception as e:
        return f"❌ Error getting file info for {filename}: {str(e)}"


def get_file_tools(executor_work_dir: Optional[str] = None) -> List[FunctionTool]:
    """
    Get universal file tools for any agent.
    
    Args:
        executor_work_dir: Working directory for file operations
        
    Returns:
        List of file management tools
    """
    tools = []
    
    # Create partial functions with work_dir bound
    # Functions return strings - FunctionTool will handle formatting
    def bound_list_files_typed() -> str:
        return asyncio.run(list_files_in_work_dir(executor_work_dir))
    
    def bound_read_file_typed(filename: str, max_length: int = 5000) -> str:
        return asyncio.run(read_file_from_work_dir(filename, executor_work_dir, max_length))
    
    def bound_get_file_info_typed(filename: str) -> str:
        return asyncio.run(get_file_info(filename, executor_work_dir))
    
    def bound_create_file_typed(filename: str, content: str) -> str:
        return asyncio.run(create_file(filename, content, executor_work_dir))
    
    # Import create_safe_function_tool locally to avoid circular import
    from agents import create_safe_function_tool
    
    # Add tools using create_safe_function_tool to ensure proper JSON wrapping
    tools.append(create_safe_function_tool(
        bound_list_files_typed,
        description="Intelligently discover and categorize all files in the working directory with comprehensive metadata"
    ))
    
    tools.append(create_safe_function_tool(
        bound_read_file_typed,
        description="Intelligently read and analyze any file type with automatic content detection and preview generation"
    ))
    
    tools.append(create_safe_function_tool(
        bound_get_file_info_typed,
        description="Get comprehensive metadata and analysis for any file type including permissions and recommendations"
    ))
    
    tools.append(create_safe_function_tool(
        bound_create_file_typed,
        description="Create a new file with the given content in the working directory. Use this to save JSON data files, CSV files, or any text-based files."
    ))

    # Add a convenience uploader for the newest deliverables
    async def _upload_recent_deliverables(max_age_minutes: int = 15, max_files: int = 5) -> str:
        try:
            from .azure_blob_tools import upload_file_to_azure_blob
            import time
            work_dir = executor_work_dir or os.getcwd()
            now = time.time()
            deliverable_exts = {".pptx", ".ppt", ".csv", ".png", ".jpg", ".jpeg", ".pdf", ".zip"}
            candidates: List[str] = []
            if os.path.isdir(work_dir):
                for name in os.listdir(work_dir):
                    path = os.path.join(work_dir, name)
                    if os.path.isfile(path) and os.path.splitext(name)[1].lower() in deliverable_exts:
                        try:
                            mtime = os.path.getmtime(path)
                            if now - mtime <= max_age_minutes * 60:
                                candidates.append(path)
                        except Exception:
                            continue
            candidates.sort(key=lambda p: os.path.getmtime(p), reverse=True)
            uploads = []
            for p in candidates[:max_files]:
                try:
                    up_json = upload_file_to_azure_blob(p)
                    uploads.append(json.loads(up_json))
                except Exception as e:
                    uploads.append({"error": str(e), "file": p})
            return json.dumps({"uploads": uploads})
        except Exception as e:
            return json.dumps({"error": str(e)})

    def bound_upload_recent_deliverables_typed(max_age_minutes: int = 15, max_files: int = 5) -> str:
        return asyncio.run(_upload_recent_deliverables(max_age_minutes, max_files))

    tools.append(create_safe_function_tool(
        bound_upload_recent_deliverables_typed,
        description="Upload the newest deliverables from the working directory (scans last N minutes) and return their URLs"
    ))

    # A suggestion-only tool: list likely deliverables without uploading
    async def _list_recent_deliverables(max_age_minutes: int = 15, max_files: int = 10, min_size_bytes: int = 1024) -> str:
        try:
            import time
            work_dir = executor_work_dir or os.getcwd()
            now = time.time()
            deliverable_exts = {".pptx", ".ppt", ".csv", ".png", ".jpg", ".jpeg", ".pdf", ".zip"}
            suggestions = []
            if os.path.isdir(work_dir):
                for name in os.listdir(work_dir):
                    path = os.path.join(work_dir, name)
                    if not os.path.isfile(path):
                        continue
                    ext = os.path.splitext(name)[1].lower()
                    if ext not in deliverable_exts:
                        continue
                    try:
                        size = os.path.getsize(path)
                        if size < min_size_bytes:
                            continue
                        mtime = os.path.getmtime(path)
                        age_s = now - mtime
                        if age_s > max_age_minutes * 60:
                            continue
                        suggestions.append({
                            "filename": name,
                            "absolute_path": path,
                            "size_bytes": size,
                            "age_seconds": int(age_s),
                            "extension": ext,
                        })
                    except Exception:
                        continue
            # Sort by size desc then recency
            suggestions.sort(key=lambda x: (x["size_bytes"], -x["age_seconds"]), reverse=True)
            return json.dumps({"suggestions": suggestions[:max_files]})
        except Exception as e:
            return json.dumps({"error": str(e)})

    def bound_list_recent_deliverables_typed(max_age_minutes: int = 15, max_files: int = 10, min_size_bytes: int = 1024) -> str:
        return asyncio.run(_list_recent_deliverables(max_age_minutes, max_files, min_size_bytes))

    tools.append(create_safe_function_tool(
        bound_list_recent_deliverables_typed,
        description="List likely deliverables (by type, size, recency) without uploading; returns suggestions for human-like selection"
    ))
    
    logger.info(f"✅ Universal file tools created for work_dir: {executor_work_dir}")
    return tools


# =============================================================================
# FILE PREVIEW GENERATION FUNCTIONS
# =============================================================================

def generate_file_preview(file_path: str, work_dir: str) -> Optional[str]:
    """
    Generate a preview image/thumbnail for supported file types.

    Args:
        file_path: Path to the file to preview
        work_dir: Working directory for saving preview images

    Returns:
        Path to generated preview image, or None if preview cannot be generated
    """
    logger.info(f"🎨 generate_file_preview called for {file_path}")
    if not os.path.exists(file_path):
        logger.warning(f"🎨 File does not exist: {file_path}")
        return None

    filename = os.path.basename(file_path)
    extension = os.path.splitext(filename)[1].lower()
    logger.info(f"🎨 Processing {filename} with extension {extension}")

    try:
        if extension == '.pdf':
            logger.info(f"🎨 Generating PDF preview for {filename}")
            return _generate_pdf_preview(file_path, work_dir)
        elif extension == '.pptx':
            logger.info(f"🎨 Generating PPTX preview for {filename}")
            return _generate_pptx_preview(file_path, work_dir)
        elif extension in ['.docx', '.xlsx']:
            logger.info(f"🎨 Generating Office preview for {filename}")
            return _generate_office_preview(file_path, work_dir)
        else:
            logger.info(f"🎨 No preview support for extension {extension}")
            return None
    except Exception as e:
        logger.warning(f"Failed to generate preview for {filename}: {e}")
        return None


def _generate_pdf_preview(file_path: str, work_dir: str) -> Optional[str]:
    """Generate ONE comprehensive preview image for PDF files - combines multiple pages."""
    try:
        from pdf2image import convert_from_path
        from PIL import Image
        import tempfile

        # Get PDF info to determine number of pages
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        total_pages = len(reader.pages)

        # Convert up to first 4 pages (or all pages if fewer than 4)
        max_pages = min(4, total_pages)
        images = convert_from_path(file_path, first_page=1, last_page=max_pages, dpi=100)
        if not images:
            return None

        if len(images) == 1:
            # Single page PDF - just use the first page
            preview_image = images[0]
        else:
            # Multi-page PDF - create a contact sheet combining pages
            # Calculate grid layout (2x2 for up to 4 pages)
            cols = 2
            rows = (len(images) + 1) // 2  # Ceiling division

            # Thumbnail size for each page
            thumb_width = 300
            thumb_height = 400

            # Create combined image
            combined_width = cols * thumb_width
            combined_height = rows * thumb_height
            combined_image = Image.new('RGB', (combined_width, combined_height), color='white')

            # Add each page as a thumbnail
            for i, page_img in enumerate(images):
                # Resize page to thumbnail
                page_img.thumbnail((thumb_width, thumb_height), Image.Resampling.LANCZOS)

                # Calculate position in grid
                x = (i % cols) * thumb_width
                y = (i // cols) * thumb_height

                # Center the thumbnail in its cell
                cell_center_x = x + thumb_width // 2
                cell_center_y = y + thumb_height // 2
                img_center_x = page_img.width // 2
                img_center_y = page_img.height // 2

                paste_x = cell_center_x - img_center_x
                paste_y = cell_center_y - img_center_y

                combined_image.paste(page_img, (paste_x, paste_y))

                # Add page number label
                from PIL import ImageDraw, ImageFont
                try:
                    draw = ImageDraw.Draw(combined_image)
                    # Try to use a small font
                    try:
                        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
                    except:
                        font = ImageFont.load_default()

                    # Add page number in bottom-right of each thumbnail
                    label = f"Page {i+1}"
                    bbox = draw.textbbox((paste_x, paste_y + thumb_height - 25), label, font=font)
                    draw.rectangle(bbox, fill="white")
                    draw.text((paste_x + 5, paste_y + thumb_height - 25), label, fill="black", font=font)
                except Exception as e:
                    logger.debug(f"Could not add page labels to PDF preview: {e}")

            preview_image = combined_image

        # Save the comprehensive preview image
        filename = os.path.basename(file_path)
        preview_filename = f"preview_{os.path.splitext(filename)[0]}.png"
        preview_path = os.path.join(work_dir, preview_filename)

        preview_image.save(preview_path, 'PNG')
        logger.info(f"🎨 Generated comprehensive PDF preview with {len(images)} pages: {preview_path}")
        return preview_path

    except ImportError as e:
        logger.warning(f"Missing dependencies for PDF preview generation: {e}")
        return None
    except Exception as e:
        logger.warning(f"PDF preview generation failed: {e}")
        return None


def _generate_pptx_preview(file_path: str, work_dir: str) -> Optional[str]:
    """Generate preview image for PowerPoint files by rendering actual slides."""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont

        # Load the presentation
        prs = Presentation(file_path)
        if len(prs.slides) == 0:
            logger.warning("PowerPoint file has no slides")
            return None

        # Get filename early for use throughout function
        filename = os.path.basename(file_path)
        
        # Render first slide (title slide) as preview
        slide = prs.slides[0]
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert EMU to pixels (assuming 96 DPI)
        width_px = int(slide_width / 914400 * 96)
        height_px = int(slide_height / 914400 * 96)
        
        # Create PIL image with white background
        img = Image.new('RGB', (width_px, height_px), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to render text from shapes
        try:
            # Default font
            try:
                font_large = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
                font_medium = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
                font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
            except:
                font_large = ImageFont.load_default()
                font_medium = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            y_pos = 100
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Determine font size based on shape type
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if shape.placeholder_format.idx == 0:  # Title placeholder
                            font = font_large
                        else:
                            font = font_medium
                    else:
                        font = font_medium
                    
                    # Wrap text if too long
                    max_width = width_px - 200
                    words = text.split()
                    lines = []
                    current_line = []
                    for word in words:
                        test_line = ' '.join(current_line + [word])
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        if bbox[2] - bbox[0] <= max_width:
                            current_line.append(word)
                        else:
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    # Draw text lines
                    for line in lines[:5]:  # Limit to 5 lines
                        if y_pos + 60 > height_px - 100:
                            break
                        draw.text((100, y_pos), line, fill='black', font=font)
                        y_pos += 60
                    
                    if y_pos > height_px - 100:
                        break
            
            # If no text found, add filename
            if y_pos == 100:
                draw.text((100, height_px // 2), filename, fill='black', font=font_large)
        except Exception as e:
            logger.warning(f"Error rendering slide text: {e}")
            # Fallback: just show filename
            try:
                font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
            except:
                font = ImageFont.load_default()
            draw.text((100, height_px // 2), filename, fill='black', font=font)

        # Save preview image
        preview_filename = f"preview_{os.path.splitext(filename)[0]}.png"
        preview_path = os.path.join(work_dir, preview_filename)
        
        img.save(preview_path, 'PNG', dpi=(150, 150))
        logger.info(f"✅ Generated PPTX preview: {preview_path}")
        return preview_path

    except ImportError:
        logger.warning("python-pptx or PIL not available for PPTX preview generation")
        return None
    except Exception as e:
        logger.warning(f"PowerPoint preview generation failed: {e}")
        return None


def _generate_office_preview(file_path: str, work_dir: str) -> Optional[str]:
    """Generate preview image for Word/Excel files."""
    try:
        import matplotlib.pyplot as plt
        import pandas as pd
        from matplotlib.patches import Rectangle

        filename = os.path.basename(file_path)
        extension = os.path.splitext(filename)[1].lower()

        if extension == '.xlsx':
            # For Excel files, try to read and display actual content
            try:
                # Read the first sheet
                df = pd.read_excel(file_path, nrows=20)  # Limit to first 20 rows for preview
                
                # Create figure with appropriate size
                fig, ax = plt.subplots(figsize=(12, min(8, max(4, len(df) * 0.3 + 2))))
                ax.set_facecolor('#ffffff')
                ax.axis('off')
                
                # Create table
                if len(df) > 0:
                    # Prepare data for table (limit columns for readability)
                    display_df = df.head(15)  # Show first 15 rows
                    if len(display_df.columns) > 8:
                        display_df = display_df.iloc[:, :8]  # Limit to 8 columns
                    
                    # Create table
                    table = ax.table(
                        cellText=display_df.values.tolist(),
                        colLabels=[str(col)[:20] for col in display_df.columns],  # Truncate long column names
                        cellLoc='left',
                        loc='center',
                        bbox=[0, 0, 1, 1]
                    )
                    
                    # Style the table
                    table.auto_set_font_size(False)
                    table.set_fontsize(8)
                    table.scale(1, 1.5)
                    
                    # Style header row
                    for i in range(len(display_df.columns)):
                        cell = table[(0, i)]
                        cell.set_facecolor('#2563eb')
                        cell.set_text_props(weight='bold', color='white')
                    
                    # Style data rows (alternating colors)
                    for i in range(1, len(display_df) + 1):
                        for j in range(len(display_df.columns)):
                            cell = table[(i, j)]
                            if i % 2 == 0:
                                cell.set_facecolor('#f3f4f6')
                            else:
                                cell.set_facecolor('#ffffff')
                            cell.set_text_props(color='#000000')
                    
                    # Add title
                    ax.text(0.5, 0.98, f"Excel Workbook: {filename}", 
                           ha='center', va='top', fontsize=12, fontweight='bold',
                           transform=ax.transAxes)
                    
                    if len(df) > 15:
                        ax.text(0.5, 0.02, f"Showing first 15 of {len(df)} rows", 
                               ha='center', va='bottom', fontsize=8, style='italic',
                               transform=ax.transAxes)
                else:
                    # Empty Excel file
                    ax.text(0.5, 0.5, f"Excel Workbook: {filename}\n(Empty file)", 
                           ha='center', va='center', fontsize=12)
                
            except Exception as e:
                logger.warning(f"Failed to read Excel content for preview: {e}, using fallback")
                # Fallback to simple preview
                ax.text(0.5, 0.7, f"Excel Workbook", ha='center', va='center',
                        fontsize=16, fontweight='bold')
                ax.text(0.5, 0.5, f"{filename}", ha='center', va='center',
                        fontsize=12)
                ax.text(0.5, 0.3, "Click download link below to view", ha='center', va='center',
                        fontsize=10, style='italic')
                ax.set_xlim(0, 1)
                ax.set_ylim(0, 1)
        elif extension == '.docx':
            # Word documents - keep simple preview for now
            ax.text(0.5, 0.7, f"Word Document", ha='center', va='center',
                    fontsize=16, fontweight='bold')
            ax.text(0.5, 0.5, f"{filename}", ha='center', va='center',
                    fontsize=12)
            ax.text(0.5, 0.3, "Click download link below to view", ha='center', va='center',
                    fontsize=10, style='italic')
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis('off')

        # Save preview image
        preview_filename = f"preview_{os.path.splitext(filename)[0]}.png"
        preview_path = os.path.join(work_dir, preview_filename)

        plt.savefig(preview_path, dpi=150, bbox_inches='tight', facecolor='#ffffff')
        plt.close()

        return preview_path

    except Exception as e:
        logger.warning(f"Office document preview generation failed: {e}")
        return None


def extract_pdf_text(file_path: str) -> str:
    """Extract text content from PDF files to validate content quality. Returns JSON with extracted text and validation."""
    try:
        import json
        import os

        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}", "is_valid": False})

        filename = os.path.basename(file_path)

        # Extract text using pypdf
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        extracted_text = ""

        # Extract text from first few pages (limit to avoid huge content)
        max_pages = min(5, len(reader.pages))
        for page_num in range(max_pages):
            page = reader.pages[page_num]
            text = page.extract_text()
            if text.strip():
                extracted_text += f"\n--- Page {page_num + 1} ---\n{text.strip()}"

        # Validate content
        is_valid = True
        validation_errors = []

        if not extracted_text.strip():
            is_valid = False
            validation_errors.append("PDF contains no extractable text content")

        # Check for error messages
        error_indicators = [
            "generation failed",
            "contact admin",
            "error",
            "failed",
            "unable to",
            "cannot",
            "report generation failed"
        ]

        lower_text = extracted_text.lower()
        for indicator in error_indicators:
            if indicator in lower_text:
                is_valid = False
                validation_errors.append(f"PDF contains error message: '{indicator}'")

        # Check for minimum content length
        if len(extracted_text.strip()) < 100:
            is_valid = False
            validation_errors.append("PDF content too short (less than 100 characters)")

        return json.dumps({
            "filename": filename,
            "extracted_text": extracted_text[:2000],  # Limit for LLM context
            "total_pages": len(reader.pages),
            "is_valid": is_valid,
            "validation_errors": validation_errors,
            "content_length": len(extracted_text)
        })

    except Exception as e:
        logger.error(f"PDF text extraction failed for {file_path}: {e}")
        return json.dumps({
            "error": str(e),
            "is_valid": False,
            "validation_errors": ["PDF text extraction failed"]
        })


def extract_pdf_text(file_path: str) -> str:
    """Extract text content from PDF files to validate content quality. Returns JSON with extracted text and validation."""
    try:
        import json
        import os

        if not os.path.exists(file_path):
            return json.dumps({"error": f"File not found: {file_path}", "is_valid": False})

        filename = os.path.basename(file_path)

        # Extract text using pypdf
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        extracted_text = ""

        # Extract text from first few pages (limit to avoid huge content)
        max_pages = min(5, len(reader.pages))
        for page_num in range(max_pages):
            page = reader.pages[page_num]
            text = page.extract_text()
            if text.strip():
                extracted_text += f"\n--- Page {page_num + 1} ---\n{text.strip()}"

        # Validate content
        is_valid = True
        validation_errors = []

        if not extracted_text.strip():
            is_valid = False
            validation_errors.append("PDF contains no extractable text content")

        # Check for error messages
        error_indicators = [
            "generation failed",
            "contact admin",
            "error",
            "failed",
            "unable to",
            "cannot",
            "report generation failed"
        ]

        lower_text = extracted_text.lower()
        for indicator in error_indicators:
            if indicator in lower_text:
                is_valid = False
                validation_errors.append(f"PDF contains error message: '{indicator}'")

        # Check for minimum content length
        if len(extracted_text.strip()) < 100:
            is_valid = False
            validation_errors.append("PDF content too short (less than 100 characters)")

        return json.dumps({
            "filename": filename,
            "extracted_text": extracted_text[:2000],  # Limit for LLM context
            "total_pages": len(reader.pages),
            "is_valid": is_valid,
            "validation_errors": validation_errors,
            "content_length": len(extracted_text)
        })

    except Exception as e:
        logger.error(f"PDF text extraction failed for {file_path}: {e}")
        return json.dumps({
            "error": str(e),
            "is_valid": False,
            "validation_errors": ["PDF text extraction failed"]
        })

def extract_pptx_text(file_path: str) -> str:
    """Extract text content from PowerPoint files to validate content quality. Returns JSON with extracted text and validation."""
    try:
        from pptx import Presentation
        import json
        
        presentation = Presentation(file_path)
        extracted_text = ""
        slide_count = 0
        total_shapes = 0
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_count += 1
            slide_text = ""
            shapes_with_text = 0
            
            for shape in slide.shapes:
                total_shapes += 1
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_with_text += 1
                    slide_text += shape.text.strip() + " "
            
            if slide_text.strip():
                extracted_text += f"\n--- Slide {slide_num + 1} ---\n{slide_text.strip()}"
        
        # Validation logic
        validation_errors = []
        
        # Check if presentation has content
        if slide_count == 0:
            validation_errors.append("PPTX contains no slides")
        elif not extracted_text.strip():
            validation_errors.append("PPTX contains no extractable text content")
        
        # Check for error messages in content
        lower_text = extracted_text.lower()
        error_indicators = [
            'error: unable to generate',
            'generation failed', 
            'contact admin',
            'system error',
            'unable to create',
            'failed to generate',
            'character at index',
            'outside the range of characters supported by the font',
            'font error',
            'unable to render',
            'presentation creation failed'
        ]
        
        for error_msg in error_indicators:
            if error_msg in lower_text:
                validation_errors.append(f"Contains error message: '{error_msg}'")
        
        # Check minimum content length
        if len(extracted_text.strip()) < 50:
            validation_errors.append("PPTX content is too short (less than 50 characters)")
        
        return json.dumps({
            "is_valid": len(validation_errors) == 0,
            "validation_errors": validation_errors,
            "text": extracted_text,
            "slide_count": slide_count,
            "total_shapes": total_shapes,
            "content_length": len(extracted_text)
        })
        
    except ImportError:
        return json.dumps({
            "error": "python-pptx library not available",
            "is_valid": False,
            "validation_errors": ["Missing python-pptx library for PPTX text extraction"]
        })
    except Exception as e:
        logger.error(f"PPTX text extraction failed for {file_path}: {e}")
        return json.dumps({
            "error": str(e),
            "is_valid": False,
            "validation_errors": ["PPTX text extraction failed"]
        })
